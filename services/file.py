import os
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx

import chardet # to support non utf-8 encoding
import tempfile # to support windows (does not have /tmp/ folder like linux)

from models.models import Document, DocumentMetadata


async def get_document_from_file(
    file: UploadFile, metadata: DocumentMetadata
) -> Document:
    extracted_text = await extract_text_from_form_file(file)

    doc = Document(text=extracted_text, metadata=metadata)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    try:
        with open(filepath, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        print(f"Error: {e}")
        raise e

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        try:
            # Try to decode the file using UTF-8
            extracted_text = file.read().decode("utf-8")
        except UnicodeDecodeError:
            # If decoding using UTF-8 fails, use chardet to detect the encoding
            file.seek(0)  # Reset file pointer to the beginning
            file_content = file.read()
            detected_encoding = chardet.detect(file_content)["encoding"]
            # Decode the file using the detected encoding
            extracted_text = file_content.decode(detected_encoding)
    elif mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text


# # Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile) -> str:
    """Return the text content of a file."""
    # Get the file body from the upload file object
    mimetype = file.content_type
    print(f"mimetype: {mimetype}")
    print(f"file.file: {file.file}")
    print("file: ", file)

    file_stream = await file.read()

    # Use NamedTemporaryFile to create a temporary file
    with tempfile.NamedTemporaryFile(mode="wb", delete=False) as f:
        # Write the contents of the uploaded file to the temporary file
        f.write(file_stream)
        # Get the path of the temporary file
        temp_file_path = f.name

    try:
        # Perform text extraction using the temporary file path
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        print(f"Error: {e}")
        os.unlink(temp_file_path)  # Remove file from temp location
        raise e

    # Remove file from temp location
    os.unlink(temp_file_path)

    return extracted_text
